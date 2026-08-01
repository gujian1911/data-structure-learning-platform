# -*- coding: utf-8 -*-
"""修复PPTX中重叠的文本框，然后导出为高清PNG
   策略：
   1. PPT -> PPTX (WPS COM)
   2. 用 python-pptx 检测并修复文本框重叠（将下方文本框下移）
   3. PPTX -> PDF -> PNG (WPS COM + PyMuPDF)
"""
import os
import sys
import time
from pptx import Presentation
from pptx.util import Emu, Inches

PPT_PATH = r"c:\Users\xuhui\Desktop\基于AIGC知识库的数据结构课程学习平台\ppt\第二章-线性表.ppt"
OUT_DIR  = r"c:\Users\xuhui\Desktop\基于AIGC知识库的数据结构课程学习平台\slides\ch2"
PPTX_PATH = os.path.join(OUT_DIR, "_fixed.pptx")
PDF_PATH  = os.path.join(OUT_DIR, "_temp.pdf")
ZOOM = 3.0

os.makedirs(OUT_DIR, exist_ok=True)

# 清理旧幻灯片
for f in os.listdir(OUT_DIR):
    if f.startswith("slide_") and f.endswith(".png"):
        os.remove(os.path.join(OUT_DIR, f))

# ============ 第1步：PPT -> PPTX ============
print("[1/5] 转换 PPT -> PPTX...")
import win32com.client
import pythoncom

pythoncom.CoInitialize()
ppt_app = win32com.client.Dispatch("PowerPoint.Application")
ppt_app.Visible = 1
pres = ppt_app.Presentations.Open(PPT_PATH, WithWindow=False)
pres.SaveAs(PPTX_PATH, 24)
pres.Close()
ppt_app.Quit()
pythoncom.CoUninitialize()
time.sleep(1)
print("      PPTX 已生成")

# ============ 第2步：修复重叠文本框 ============
print("[2/5] 修复重叠文本框...")
prs = Presentation(PPTX_PATH)
slide_h = prs.slide_height
slide_h_in = Emu(slide_h).inches

def get_rect(shape):
    """返回 (left, top, width, height) in EMU"""
    return (shape.left, shape.top, shape.width, shape.height)

def rects_overlap(r1, r2):
    l1, t1, w1, h1 = r1
    l2, t2, w2, h2 = r2
    if any(v is None for v in [l1, t1, w1, h1, l2, t2, w2, h2]):
        return False
    return not (l1 + w1 <= l2 or l2 + w2 <= l1 or t1 + h1 <= t2 or t2 + h2 <= t1)

def is_text_shape(shape):
    """判断是否为文本类形状（有文本内容的文本框/占位符/矩形）"""
    if not shape.has_text_frame:
        return False
    text = shape.text_frame.text.strip()
    return len(text) > 0

total_fixes = 0
for idx, slide in enumerate(prs.slides, 1):
    # 收集有文本的形状
    text_shapes = []
    for shape in slide.shapes:
        if shape.left is None or shape.top is None:
            continue
        if shape.width is None or shape.height is None:
            continue
        text_shapes.append(shape)

    # 按垂直位置排序
    text_shapes.sort(key=lambda s: (s.top or 0, s.left or 0))

    # 检测并修复重叠
    fixes_this_slide = 0
    changed = True
    iterations = 0
    while changed and iterations < 5:
        changed = False
        iterations += 1
        for i in range(len(text_shapes)):
            for j in range(i + 1, len(text_shapes)):
                s1 = text_shapes[i]
                s2 = text_shapes[j]
                r1 = get_rect(s1)
                r2 = get_rect(s2)

                if not rects_overlap(r1, r2):
                    continue

                # 如果水平方向不重叠（左右排列），跳过
                l1, t1, w1, h1 = r1
                l2, t2, w2, h2 = r2
                if l1 + w1 <= l2 + Emu(0.2 * 914400) or l2 + w2 <= l1 + Emu(0.2 * 914400):
                    continue

                # 如果是空占位符，跳过（不移动它）
                if not is_text_shape(s1) and not is_text_shape(s2):
                    continue

                # 如果其中一个没有文本（空占位符），尝试移除它
                if not is_text_shape(s1):
                    # s1是空的，s2有文本，尝试把s1移到画外
                    s1.top = Emu(int(10 * 914400))
                    changed = True
                    fixes_this_slide += 1
                    continue
                if not is_text_shape(s2):
                    s2.top = Emu(int(10 * 914400))
                    changed = True
                    fixes_this_slide += 1
                    continue

                # 两个都有文本，将下方的形状下移
                # s1在上，s2在下（因为已排序）
                new_top = s1.top + s1.height + Emu(int(0.08 * 914400))  # 0.08英寸间距

                # 确保不超出幻灯片底部
                if new_top + s2.height > slide_h:
                    # 如果超出，缩小高度
                    s2.height = slide_h - new_top - Emu(int(0.2 * 914400))

                s2.top = new_top
                changed = True
                fixes_this_slide += 1

    if fixes_this_slide > 0:
        print(f"      第{idx}页: 修复了 {fixes_this_slide} 处重叠")
        total_fixes += fixes_this_slide

prs.save(PPTX_PATH)
print(f"      总共修复 {total_fixes} 处重叠")

# ============ 第3步：PPTX -> PDF ============
print("[3/5] 转换 PPTX -> PDF...")
pythoncom.CoInitialize()
ppt_app = win32com.client.Dispatch("PowerPoint.Application")
ppt_app.Visible = 1
pres = ppt_app.Presentations.Open(PPTX_PATH, WithWindow=False)
pres.SaveAs(PDF_PATH, 32)  # PDF
pres.Close()
ppt_app.Quit()
pythoncom.CoUninitialize()
time.sleep(1)
print("      PDF 已生成")

# ============ 第4步：PDF -> PNG ============
print(f"[4/5] 渲染 PDF -> PNG（{ZOOM}x缩放）...")
import fitz

doc = fitz.open(PDF_PATH)
total = len(doc)
matrix = fitz.Matrix(ZOOM, ZOOM)

for i in range(total):
    page = doc[i]
    pix = page.get_pixmap(matrix=matrix, alpha=False)
    out_path = os.path.join(OUT_DIR, f"slide_{str(i + 1).zfill(3)}.png")
    pix.save(out_path)
    print(f"      slide_{str(i + 1).zfill(3)}.png  ({pix.width}x{pix.height})")
    pix = None
    page = None

doc.close()

# ============ 第5步：清理 ============
print("[5/5] 清理临时文件...")
for f in [PPTX_PATH, PDF_PATH]:
    if os.path.exists(f):
        os.remove(f)

print(f"\n完成！共 {total} 张幻灯片已生成（修复了 {total_fixes} 处文本框重叠）。")
