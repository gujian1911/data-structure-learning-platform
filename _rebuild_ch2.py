# -*- coding: utf-8 -*-
"""一次性完成：PPT -> PPTX -> 修复重叠 -> PDF -> PNG"""
import os, sys, time
import subprocess

OUT_DIR  = r"c:\Users\xuhui\Desktop\基于AIGC知识库的数据结构课程学习平台\slides\ch2"
PPTX_PATH = os.path.join(OUT_DIR, "_fixed.pptx")
PDF_PATH  = os.path.join(OUT_DIR, "_temp.pdf")

# 清理旧幻灯片
for f in os.listdir(OUT_DIR):
    if f.startswith("slide_") and f.endswith(".png"):
        os.remove(os.path.join(OUT_DIR, f))
print("[0] 已清理旧幻灯片")

# ===== 步骤1：PPT -> PPTX =====
print("[1] 转换 PPT -> PPTX...")
import win32com.client
import pythoncom

pythoncom.CoInitialize()
ppt_app = win32com.client.Dispatch("PowerPoint.Application")
ppt_app.Visible = 1
pres = ppt_app.Presentations.Open(
    r"c:\Users\xuhui\Desktop\基于AIGC知识库的数据结构课程学习平台\ppt\第二章-线性表.ppt",
    WithWindow=False
)
pres.SaveAs(PPTX_PATH, 24)
pres.Close()
ppt_app.Quit()
pythoncom.CoUninitialize()
time.sleep(2)
print("    PPTX 已生成")

# ===== 步骤2：修复重叠文本框 =====
print("[2] 修复重叠文本框...")
from pptx import Presentation
from pptx.util import Emu

prs = Presentation(PPTX_PATH)
slide_h = prs.slide_height

def rects_overlap(s1, s2):
    l1, t1, w1, h1 = s1.left, s1.top, s1.width, s1.height
    l2, t2, w2, h2 = s2.left, s2.top, s2.width, s2.height
    if any(v is None for v in [l1, t1, w1, h1, l2, t2, w2, h2]):
        return False
    return not (l1 + w1 <= l2 or l2 + w2 <= l1 or t1 + h1 <= t2 or t2 + h2 <= t1)

def has_text(shape):
    if not shape.has_text_frame:
        return False
    return len(shape.text_frame.text.strip()) > 0

GAP = Emu(int(0.08 * 914400))
HORIZ_TOL = Emu(int(0.3 * 914400))
total_fixes = 0

for idx, slide in enumerate(prs.slides, 1):
    shapes = [s for s in slide.shapes
              if s.left is not None and s.top is not None
              and s.width is not None and s.height is not None]
    shapes.sort(key=lambda s: (s.top, s.left))

    fixes = 0
    changed = True
    for _ in range(5):
        if not changed:
            break
        changed = False
        for i in range(len(shapes)):
            for j in range(i + 1, len(shapes)):
                s1, s2 = shapes[i], shapes[j]
                if not rects_overlap(s1, s2):
                    continue
                if s1.left + s1.width <= s2.left + HORIZ_TOL or s2.left + s2.width <= s1.left + HORIZ_TOL:
                    continue
                t1, t2 = has_text(s1), has_text(s2)
                if not t1 and not t2:
                    continue
                if not t1:
                    s1.top = Emu(int(10 * 914400))
                    changed = True; fixes += 1; continue
                if not t2:
                    s2.top = Emu(int(10 * 914400))
                    changed = True; fixes += 1; continue
                new_top = s1.top + s1.height + GAP
                if new_top + s2.height > slide_h:
                    s2.height = max(slide_h - new_top - Emu(int(0.2 * 914400)), Emu(int(0.5 * 914400)))
                s2.top = new_top
                changed = True; fixes += 1

    if fixes > 0:
        print(f"    第{idx}页: 修复 {fixes} 处")
        total_fixes += fixes

prs.save(PPTX_PATH)
print(f"    总共修复 {total_fixes} 处重叠")

# ===== 步骤3：PPTX -> PDF =====
print("[3] 转换 PPTX -> PDF...")
pythoncom.CoInitialize()
ppt_app = win32com.client.Dispatch("PowerPoint.Application")
ppt_app.Visible = 1
pres = ppt_app.Presentations.Open(PPTX_PATH, WithWindow=False)
pres.SaveAs(PDF_PATH, 32)
pres.Close()
ppt_app.Quit()
pythoncom.CoUninitialize()
time.sleep(2)

# 验证PDF是否存在
if not os.path.exists(PDF_PATH):
    print("[错误] PDF 生成失败！")
    sys.exit(1)
print("    PDF 已生成")

# ===== 步骤4：PDF -> PNG =====
ZOOM = 3.0
print(f"[4] 渲染 PDF -> PNG（{ZOOM}x缩放）...")
import fitz

doc = fitz.open(PDF_PATH)
total = len(doc)
matrix = fitz.Matrix(ZOOM, ZOOM)

for i in range(total):
    page = doc[i]
    pix = page.get_pixmap(matrix=matrix, alpha=False)
    out_path = os.path.join(OUT_DIR, f"slide_{str(i + 1).zfill(3)}.png")
    pix.save(out_path)
    print(f"    slide_{str(i + 1).zfill(3)}.png  ({pix.width}x{pix.height})")
    pix = None
    page = None

doc.close()

# ===== 清理 =====
for f in [PPTX_PATH, PDF_PATH]:
    if os.path.exists(f):
        os.remove(f)

print(f"\n完成！共 {total} 张幻灯片已生成（修复了 {total_fixes} 处文本框重叠）。")
