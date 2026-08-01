# -*- coding: utf-8 -*-
"""检查PPTX中每张幻灯片的形状布局，检测重叠的文本框"""
import os
import sys
import time

PPT_PATH = r"c:\Users\xuhui\Desktop\基于AIGC知识库的数据结构课程学习平台\ppt\第二章-线性表.ppt"
PPTX_PATH = r"c:\Users\xuhui\Desktop\基于AIGC知识库的数据结构课程学习平台\slides\ch2\_inspect.pptx"

# 第1步：PPT -> PPTX
print("[1/2] 转换 PPT -> PPTX...")
import win32com.client
import pythoncom

pythoncom.CoInitialize()
ppt_app = win32com.client.Dispatch("PowerPoint.Application")
ppt_app.Visible = 1
pres = ppt_app.Presentations.Open(PPT_PATH, WithWindow=False)
pres.SaveAs(PPTX_PATH, 24)  # PPTX
pres.Close()
ppt_app.Quit()
pythoncom.CoUninitialize()
time.sleep(1)
print("      PPTX 已生成")

# 第2步：检查形状布局
print("[2/2] 检查形状布局...")
from pptx import Presentation
from pptx.util import Emu, Inches, Pt

prs = Presentation(PPTX_PATH)
slide_w = prs.slide_width
slide_h = prs.slide_height
print(f"      幻灯片尺寸: {slide_w} x {slide_h} EMU ({Emu(slide_w).inches:.1f} x {Emu(slide_h).inches:.1f} in)")
print()

def rects_overlap(r1, r2):
    """检查两个矩形是否重叠"""
    l1, t1, w1, h1 = r1
    l2, t2, w2, h2 = r2
    # 重叠条件
    return not (l1 + w1 <= l2 or l2 + w2 <= l1 or t1 + h1 <= t2 or t2 + h2 <= t1)

for idx, slide in enumerate(prs.slides, 1):
    shapes_info = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text[:60].replace('\n', ' | ')
        else:
            text = f"[{shape.shape_type}]"
        info = {
            'name': shape.name,
            'left': shape.left,
            'top': shape.top,
            'width': shape.width,
            'height': shape.height,
            'text': text,
            'shape': shape,
        }
        shapes_info.append(info)

    # 检测重叠
    overlaps = []
    for i in range(len(shapes_info)):
        for j in range(i + 1, len(shapes_info)):
            r1 = (shapes_info[i]['left'], shapes_info[i]['top'], shapes_info[i]['width'], shapes_info[i]['height'])
            r2 = (shapes_info[j]['left'], shapes_info[j]['top'], shapes_info[j]['width'], shapes_info[j]['height'])
            if r1[0] is None or r2[0] is None:
                continue
            if rects_overlap(r1, r2):
                overlaps.append((i, j))

    if overlaps:
        print(f"=== 第{idx}页（{len(shapes_info)}个形状，{len(overlaps)}处重叠）===")
        for si, info in enumerate(shapes_info):
            l = Emu(info['left']).inches if info['left'] else 0
            t = Emu(info['top']).inches if info['top'] else 0
            w = Emu(info['width']).inches if info['width'] else 0
            h = Emu(info['height']).inches if info['height'] else 0
            print(f"  [{si}] {info['name']}: ({l:.2f}, {t:.2f}) {w:.2f}x{h:.2f}  \"{info['text']}\"")
        for i, j in overlaps:
            print(f"  >>> 重叠: [{i}] <-> [{j}]")
        print()

# 清理
if os.path.exists(PPTX_PATH):
    os.remove(PPTX_PATH)
