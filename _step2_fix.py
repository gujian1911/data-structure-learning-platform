# -*- coding: utf-8 -*-
"""步骤2：修复PPTX中重叠的文本框"""
import os
from pptx import Presentation
from pptx.util import Emu

PPTX_PATH = r"c:\Users\xuhui\Desktop\基于AIGC知识库的数据结构课程学习平台\slides\ch2\_fixed.pptx"

print("[2] 修复重叠文本框...")
prs = Presentation(PPTX_PATH)
slide_h = prs.slide_height

def rects_overlap(s1, s2):
    l1, t1, w1, h1 = s1.left, s1.top, s1.width, s1.height
    l2, t2, w2, h2 = s2.left, s2.top, s2.width, s2.height
    if any(v is None for v in [l1, t1, w1, h1, l2, t2, w2, h2]):
        return False
    return not (l1 + w1 <= l2 or l2 + w2 <= l1 or t1 + h1 <= t2 or t2 + h2 <= t1)

def has_text(shape):
    if not shape.has_text_frame:
        return False
    return len(shape.text_frame.text.strip()) > 0

GAP = Emu(int(0.08 * 914400))  # 0.08英寸间距
HORIZ_TOL = Emu(int(0.3 * 914400))  # 水平方向容忍度

total_fixes = 0
for idx, slide in enumerate(prs.slides, 1):
    shapes = [s for s in slide.shapes
              if s.left is not None and s.top is not None
              and s.width is not None and s.height is not None]
    shapes.sort(key=lambda s: (s.top, s.left))

    fixes = 0
    changed = True
    for _ in range(5):
        if not changed:
            break
        changed = False
        for i in range(len(shapes)):
            for j in range(i + 1, len(shapes)):
                s1, s2 = shapes[i], shapes[j]
                if not rects_overlap(s1, s2):
                    continue

                # 水平方向不重叠（左右排列）则跳过
                if s1.left + s1.width <= s2.left + HORIZ_TOL or s2.left + s2.width <= s1.left + HORIZ_TOL:
                    continue

                t1, t2 = has_text(s1), has_text(s2)
                if not t1 and not t2:
                    continue

                # 空占位符移到画外
                if not t1:
                    s1.top = Emu(int(10 * 914400))
                    changed = True
                    fixes += 1
                    continue
                if not t2:
                    s2.top = Emu(int(10 * 914400))
                    changed = True
                    fixes += 1
                    continue

                # 两个都有文本，下移 s2
                new_top = s1.top + s1.height + GAP
                if new_top + s2.height > slide_h:
                    s2.height = max(slide_h - new_top - Emu(int(0.2 * 914400)), Emu(int(0.5 * 914400)))
                s2.top = new_top
                changed = True
                fixes += 1

    if fixes > 0:
        print(f"    第{idx}页: 修复 {fixes} 处")
        total_fixes += fixes

prs.save(PPTX_PATH)
print(f"    总共修复 {total_fixes} 处重叠")
